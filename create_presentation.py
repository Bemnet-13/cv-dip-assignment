import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RESULTS_DIR = "results"
OUTPUT_FILE = "DIP_Elementary_Methods.pptx"

# ── Color palette ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xAA, 0xCC, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT     = RGBColor(0x2E, 0x86, 0xC1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Low-level helpers ──────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, no_line=True):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_header(slide, label, height=Inches(0.85)):
    add_rect(slide, 0, 0, SLIDE_W, height, NAVY)
    add_text(slide, label, Inches(0.4), Inches(0.1), Inches(12.5), height,
             size=24, bold=True, color=WHITE)


# ── Slide builders ─────────────────────────────────────────────────────────────

def title_slide(name, course, date):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, NAVY)

    add_rect(slide, Inches(0.5), Inches(2.1), Inches(12.333), Inches(0.06), ACCENT)

    add_text(slide, "Digital Image Processing",
             Inches(0.5), Inches(0.8), Inches(12.333), Inches(1.2),
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Elementary Methods",
             Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.8),
             size=30, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    add_text(slide, f"{name}  ·  {course}  ·  {date}",
             Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.6),
             size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)


def agenda_slide(methods):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, LIGHT_GRAY)
    add_header(slide, "Overview")

    for i, method in enumerate(methods):
        y = Inches(1.05) + i * Inches(0.83)
        add_rect(slide, Inches(0.5), y + Inches(0.08),
                 Inches(0.45), Inches(0.45), NAVY)
        add_text(slide, str(i + 1), Inches(0.5), y + Inches(0.06),
                 Inches(0.45), Inches(0.45), size=16, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, method, Inches(1.1), y, Inches(11.5), Inches(0.6),
                 size=18, color=DARK_TEXT)


def method_intro_slide(number, name, formula, description_lines):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, LIGHT_GRAY)
    add_header(slide, f"{number}. {name}", height=Inches(1.0))

    # Formula box
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.75), NAVY)
    add_text(slide, f"Formula:   {formula}",
             Inches(0.7), Inches(1.28), Inches(12.0), Inches(0.62),
             size=20, bold=True, color=WHITE)

    # Description lines
    y = Inches(2.2)
    for line in description_lines:
        add_text(slide, line, Inches(0.6), y, Inches(12.1), Inches(0.65),
                 size=17, color=DARK_TEXT)
        y += Inches(0.7)


def results_slide(header_label, image_path):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_header(slide, header_label, height=Inches(0.72))

    if os.path.exists(image_path):
        img_top    = Inches(0.82)
        img_height = SLIDE_H - img_top - Inches(0.18)
        img_width  = SLIDE_W - Inches(0.6)
        slide.shapes.add_picture(image_path, Inches(0.3), img_top,
                                 img_width, img_height)
    else:
        add_text(slide, f"[Image not found: {image_path}]",
                 Inches(1), Inches(3), Inches(11), Inches(1),
                 size=16, color=RGBColor(0xCC, 0x00, 0x00))


def conclusion_slide(points):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, NAVY)

    add_text(slide, "Conclusion", Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0),
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.5), Inches(1.3), Inches(10.333), Inches(0.06), ACCENT)

    y = Inches(1.6)
    for point in points:
        add_rect(slide, Inches(0.8), y + Inches(0.14), Inches(0.18), Inches(0.18), LIGHT_BLUE)
        add_text(slide, point, Inches(1.2), y, Inches(11.5), Inches(0.65),
                 size=18, color=WHITE)
        y += Inches(0.75)


# ── Content data ───────────────────────────────────────────────────────────────

METHODS = [
    "Image Negative",
    "Gamma Encoding / Correction",
    "Logarithmic Transformation",
    "Contrast Stretching",
    "Histogram Equalization",
    "Intensity Level Slicing",
    "Bit Plane Slicing",
]

METHOD_DATA = [
    {
        "name": "Image Negative",
        "formula": "s = 255 − r",
        "desc": [
            "Inverts the intensity of every pixel by subtracting it from the maximum value (255).",
            "Useful for enhancing white or grey detail embedded in dark regions of an image.",
            "Applied identically to each channel in RGB and directly to the single channel in grayscale.",
        ],
        "results": [
            ("Image Negative — RGB",       f"{RESULTS_DIR}/image_negative_rgb.png"),
            ("Image Negative — Grayscale", f"{RESULTS_DIR}/image_negative_gray.png"),
        ],
    },
    {
        "name": "Gamma Encoding / Correction",
        "formula": "s = c · rᵞ",
        "desc": [
            "A non-linear power-law transformation controlled by the exponent γ (gamma).",
            "γ < 1  →  brightens the image (expands dark regions).",
            "γ > 1  →  darkens the image (compresses bright regions).",
            "Demo: γ = 0.4 on RGB image,  γ = 2.2 on grayscale image.",
        ],
        "results": [
            ("Gamma Correction — RGB (γ=0.4)",       f"{RESULTS_DIR}/gamma_correction_rgb.png"),
            ("Gamma Correction — Grayscale (γ=2.2)", f"{RESULTS_DIR}/gamma_correction_gray.png"),
        ],
    },
    {
        "name": "Logarithmic Transformation",
        "formula": "s = c · log(1 + r),   c = 255 / log(1 + max(r))",
        "desc": [
            "Maps a narrow range of low-intensity values to a wider range of output levels.",
            "Compresses high-intensity values, expanding dark/shadow detail.",
            "Constant c normalises the output to the full [0, 255] range.",
        ],
        "results": [
            ("Logarithmic Transform — RGB",       f"{RESULTS_DIR}/logarithmic_transform_rgb.png"),
            ("Logarithmic Transform — Grayscale", f"{RESULTS_DIR}/logarithmic_transform_gray.png"),
        ],
    },
    {
        "name": "Contrast Stretching",
        "formula": "s = (r − min) / (max − min) × 255",
        "desc": [
            "Linearly rescales pixel intensities so the darkest pixel maps to 0",
            "and the brightest pixel maps to 255.",
            "Maximises the dynamic range when the original image uses only a narrow band.",
        ],
        "results": [
            ("Contrast Stretching — RGB",       f"{RESULTS_DIR}/contrast_stretching_rgb.png"),
            ("Contrast Stretching — Grayscale", f"{RESULTS_DIR}/contrast_stretching_gray.png"),
        ],
    },
    {
        "name": "Histogram Equalization",
        "formula": "s = (L−1) · CDF(r),   CDF = cumulative distribution function",
        "desc": [
            "Redistributes intensity values so the output histogram is approximately flat.",
            "For RGB: equalization is applied only to the Y (luminance) channel in YCrCb",
            "space to avoid introducing colour distortion.",
        ],
        "results": [
            ("Histogram Equalization — RGB",       f"{RESULTS_DIR}/histogram_equalization_rgb.png"),
            ("Histogram Equalization — Grayscale", f"{RESULTS_DIR}/histogram_equalization_gray.png"),
        ],
    },
    {
        "name": "Intensity Level Slicing",
        "formula": "s = 255 if lower ≤ r ≤ upper,  else r  (or 0)",
        "desc": [
            "Highlights pixels whose intensity falls within a defined range [lower, upper].",
            "Mode 1 — Preserve BG: range pixels turn white, rest unchanged.",
            "Mode 2 — Black BG: range pixels turn white, rest turn black.",
            "Range used in demo: [100, 150].  Mask built from luminance for RGB images.",
        ],
        "results": [
            ("Intensity Level Slicing — RGB  (Preserve Background)", f"{RESULTS_DIR}/intensity_slicing_preserve_bg.png"),
            ("Intensity Level Slicing — Grayscale  (Black Background)", f"{RESULTS_DIR}/intensity_slicing_black_bg.png"),
        ],
    },
    {
        "name": "Bit Plane Slicing",
        "formula": "Plane k = (pixel  AND  2ᵏ) > 0  →  255,  else 0",
        "desc": [
            "Decomposes an 8-bit image into 8 binary planes (bit 0 = LSB, bit 7 = MSB).",
            "Higher bit planes (6–7) carry the dominant structural information.",
            "Lower bit planes (0–2) capture fine detail and noise.",
            "Visualised as a 2×4 grid showing all 8 planes simultaneously.",
        ],
        "results": [
            ("Bit Plane Slicing — Grayscale Image",       f"{RESULTS_DIR}/bit_plane_slicing_grayscale_image.png"),
            ("Bit Plane Slicing — RGB Image (→ Gray)",    f"{RESULTS_DIR}/bit_plane_slicing_rgb_to_gray.png"),
        ],
    },
]

CONCLUSION_POINTS = [
    "Implemented all 7 elementary DIP intensity-transformation methods in Python.",
    "Each method handles both RGB and grayscale images correctly.",
    "RGB Histogram Equalization uses the YCrCb colour space to preserve hue.",
    "Intensity Level Slicing for RGB builds the mask from luminance (no channel artefacts).",
    "Gamma and contrast outputs are clipped before uint8 cast, preventing overflow.",
    "Results confirm clear, visible effects on both image types for every transformation.",
]


# ── Build presentation ─────────────────────────────────────────────────────────

title_slide(
    name="Bemnet Aschalew",
    course="Digital Image Processing",
    date="2026",
)

agenda_slide(METHODS)

for i, m in enumerate(METHOD_DATA):
    method_intro_slide(i + 1, m["name"], m["formula"], m["desc"])
    for label, path in m["results"]:
        results_slide(label, path)

conclusion_slide(CONCLUSION_POINTS)

prs.save(OUTPUT_FILE)
print(f"Saved: {OUTPUT_FILE}  ({len(prs.slides)} slides)")
